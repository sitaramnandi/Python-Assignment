import openpyxl
import openai
from pptx import Presentation

def read_excel_prompts(excel_file):
    prompts = []

    try:
        workbook = openpyxl.load_workbook(excel_file)
        sheet = workbook.active

        # Iterate over rows in the sheet
        for row in sheet.iter_rows(values_only=True):
            # Assuming the prompts are in the first column (column A)
            prompt = row[0]
            prompts.append(prompt)

    except FileNotFoundError:
        print(f"Excel file not found: {excel_file}")
    except Exception as e:
        print(f"Error reading Excel file: {e}")

    return prompts
def generate_short_paragraphs(prompts, api_key):
    generated_paragraphs = []

    # Set up the OpenAI API
    openai.api_key = api_key

    # Generate short paragraphs based on prompts
    for prompt in prompts:
        response = openai.Completion.create(
            engine="davinci",
            prompt=prompt,
            max_tokens=100,
            temperature=0.7,
            n=1,
            stop=None,
            frequency_penalty=0.0,
            presence_penalty=0.0
        )

        generated_paragraphs.append(response.choices[0].text.strip())

    return generated_paragraphs
def create_ppt_slides(paragraphs):
    presentation = Presentation()

    # Add a slide for each generated paragraph
    for paragraph in paragraphs:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Generated Paragraph"
        slide.placeholders[1].text = paragraph

    # Save the PowerPoint presentation
    presentation.save("output.pptx")
excel_file = "path/to/your/excel/file.xlsx"
prompts = read_excel_prompts(excel_file)

api_key = "your_openai_api_key"
generated_paragraphs = generate_short_paragraphs(prompts, api_key)

create_ppt_slides(generated_paragraphs)
